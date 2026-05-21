"""
file_service.py
---------------
Text extraction from PDF, DOCX, TXT, CSV, XLSX, PPTX.

Key fixes vs previous version:
  - PPTX: skips layout / master placeholders that cause duplicate text
  - PPTX: prefixes every slide with [Slide N] so the chunker treats
    each slide as its own section heading
  - PDF:  page_map now uses cumulative character offsets (not line offsets)
  - All formats: raw bytes are read once and rewound; UploadFile.file is
    not seekable after the first read in some FastAPI versions, so we
    read into BytesIO upfront.
  - Added near-duplicate detection at the extracted-text level:
    if a workspace already has this file's content (same SHA-256) the
    caller can skip re-embedding.
"""

from __future__ import annotations
import hashlib
import io

from fastapi import UploadFile
import pdfplumber
import docx
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import PP_PLACEHOLDER


# ── helpers ──────────────────────────────────────────────────────────────────

def _sha256(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _read_file(file: UploadFile) -> bytes:
    """Read all bytes from an UploadFile, rewinding first."""
    file.file.seek(0)
    return file.file.read()


def _is_content_placeholder(shape) -> bool:
    """
    Return True for shapes that carry the actual slide content.
    Excludes slide-layout / slide-master boilerplate placeholders
    (e.g. footer, slide number, date) that duplicate content across slides.
    """
    SKIP_TYPES = {
        PP_PLACEHOLDER.FOOTER,
        PP_PLACEHOLDER.SLIDE_NUMBER,
        PP_PLACEHOLDER.DATE,
    }
    try:
        if shape.is_placeholder and shape.placeholder_format.type in SKIP_TYPES:
            return False
    except Exception:
        pass
    return True


# ── public API ───────────────────────────────────────────────────────────────

def get_file_type(filename: str) -> str:
    filename = filename.lower()
    for ext in ["pdf", "docx", "txt", "csv", "xlsx", "xls", "pptx"]:
        if filename.endswith(f".{ext}"):
            return ext
    return "unknown"


def extract_text_from_file(file: UploadFile) -> dict:
    """
    Returns:
        {
            "text":       str,               # full extracted text
            "page_count": int | None,        # pages / slides
            "page_map":   list[tuple],       # [(char_offset, page_number), ...]
            "sha256":     str,               # content fingerprint for dedup
        }
    """
    raw = _read_file(file)
    filename = file.filename.lower()
    sha = _sha256(raw)

    # ── PDF ─────────────────────────────────────────────────────────────────
    if filename.endswith(".pdf"):
        text_parts: list[str] = []
        page_map: list[tuple[int, int]] = []
        offset = 0

        with pdfplumber.open(io.BytesIO(raw)) as pdf:
            page_count = len(pdf.pages)
            for page_num, page in enumerate(pdf.pages, start=1):
                page_text = page.extract_text() or ""
                # Record where this page starts in the concatenated text
                page_map.append((offset, page_num))
                text_parts.append(page_text)
                offset += len(page_text) + 1  # +1 for the joining newline

        full_text = "\n".join(text_parts)
        return {
            "text": full_text,
            "page_count": page_count,
            "page_map": page_map,
            "sha256": sha,
        }

    # ── DOCX ────────────────────────────────────────────────────────────────
    elif filename.endswith(".docx"):
        doc = docx.Document(io.BytesIO(raw))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return {"text": text, "page_count": None, "page_map": [], "sha256": sha}

    # ── TXT ─────────────────────────────────────────────────────────────────
    elif filename.endswith(".txt"):
        text = raw.decode("utf-8", errors="replace")
        return {"text": text, "page_count": None, "page_map": [], "sha256": sha}

    # ── CSV ─────────────────────────────────────────────────────────────────
    elif filename.endswith(".csv"):
        df = pd.read_csv(io.BytesIO(raw))
        return {"text": df.to_string(), "page_count": None, "page_map": [], "sha256": sha}

    # ── Excel ───────────────────────────────────────────────────────────────
    elif filename.endswith((".xlsx", ".xls")):
        df = pd.read_excel(io.BytesIO(raw))
        return {"text": df.to_string(), "page_count": None, "page_map": [], "sha256": sha}

    # ── PowerPoint ──────────────────────────────────────────────────────────
    elif filename.endswith(".pptx"):
        prs = Presentation(io.BytesIO(raw))
        parts: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_lines: list[str] = [f"[Slide {slide_num}]"]

            for shape in slide.shapes:
                # Skip boilerplate layout placeholders
                if not _is_content_placeholder(shape):
                    continue

                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            slide_lines.append(line)

                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                            if cell.text.strip()
                        )
                        if row_text:
                            slide_lines.append(row_text)

            # Only keep slides that have actual content beyond the marker
            if len(slide_lines) > 1:
                parts.append("\n".join(slide_lines))

        full_text = "\n\n".join(parts)
        return {
            "text": full_text,
            "page_count": len(prs.slides),
            "page_map": [],   # slide numbers are embedded in text as [Slide N]
            "sha256": sha,
        }

    # ── Unsupported ──────────────────────────────────────────────────────────
    return {
        "text": "Unsupported file format",
        "page_count": None,
        "page_map": [],
        "sha256": sha,
    }